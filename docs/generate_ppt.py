from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# Color palette
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
BLUE_ACCENT = RGBColor(0x00, 0xD4, 0xFF)
PURPLE = RGBColor(0x7B, 0x2F, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
DARK_GRAY = RGBColor(0x66, 0x66, 0x66)
TABLE_HEADER_BG = RGBColor(0x0F, 0x34, 0x60)

def set_slide_bg(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_slide(title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_bg(slide, DARK_BG)

    # Title
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.33), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = BLUE_ACCENT
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(11.33), Inches(1))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = subtitle
    p2.font.size = Pt(22)
    p2.font.color.rgb = LIGHT_GRAY
    p2.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(title, bullets=None, table_data=None, code_text=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_bg(slide, RGBColor(0x16, 0x21, 0x3E))

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = BLUE_ACCENT

    y_pos = Inches(1.3)

    if bullets:
        txBox2 = slide.shapes.add_textbox(Inches(0.8), y_pos, Inches(11.5), Inches(5.5))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf2.paragraphs[0]
            else:
                p = tf2.add_paragraph()
            p.text = bullet
            p.font.size = Pt(18)
            p.font.color.rgb = WHITE
            p.space_after = Pt(12)

    if table_data:
        rows = len(table_data)
        cols = len(table_data[0])
        tbl = slide.shapes.add_table(rows, cols, Inches(0.5), y_pos, Inches(12), Inches(rows * 0.55)).table

        for i, row_data in enumerate(table_data):
            for j, cell_text in enumerate(row_data):
                cell = tbl.cell(i, j)
                cell.text = str(cell_text)
                p = cell.text_frame.paragraphs[0]
                p.font.size = Pt(13)
                if i == 0:
                    p.font.bold = True
                    p.font.color.rgb = WHITE
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = TABLE_HEADER_BG
                else:
                    p.font.color.rgb = WHITE
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

    if code_text:
        txBox3 = slide.shapes.add_textbox(Inches(0.5), y_pos, Inches(12), Inches(5.5))
        tf3 = txBox3.text_frame
        tf3.word_wrap = False
        p = tf3.paragraphs[0]
        p.text = code_text
        p.font.size = Pt(11)
        p.font.name = "Consolas"
        p.font.color.rgb = RGBColor(0xC9, 0xD1, 0xD9)

    return slide

# ===== SLIDE 1: COVER =====
slide = add_title_slide(
    "QA Automation Framework",
    "Hiring Management System — End-to-End API Testing"
)
# Add tech badges
txBox = slide.shapes.add_textbox(Inches(2), Inches(5), Inches(9), Inches(0.6))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Java 11  |  Maven  |  TestNG 7.9  |  RestAssured 5.4  |  Gson  |  Agent-Driven"
p.font.size = Pt(16)
p.font.color.rgb = DARK_GRAY
p.alignment = PP_ALIGN.CENTER

# Add target URL
txBox2 = slide.shapes.add_textbox(Inches(2), Inches(5.8), Inches(9), Inches(0.5))
tf2 = txBox2.text_frame
p2 = tf2.paragraphs[0]
p2.text = "Target: github.com/odeshpande-conga/HiringManagementSystem"
p2.font.size = Pt(14)
p2.font.color.rgb = DARK_GRAY
p2.alignment = PP_ALIGN.CENTER

# Date
txBox3 = slide.shapes.add_textbox(Inches(2), Inches(6.5), Inches(9), Inches(0.5))
tf3 = txBox3.text_frame
p3 = tf3.paragraphs[0]
p3.text = "May 2026"
p3.font.size = Pt(14)
p3.font.color.rgb = DARK_GRAY
p3.alignment = PP_ALIGN.CENTER

# ===== SLIDE 2: AGENDA =====
add_content_slide("📋 Agenda", bullets=[
    "1.  Project Overview & Objectives",
    "2.  Technology Stack",
    "3.  Framework Architecture",
    "4.  Agent-Driven QA Pipeline (6 Steps)",
    "5.  API Coverage (16 APIs)",
    "6.  UML & Sequence Diagrams",
    "7.  Key Design Decisions",
    "8.  Scalability & Next Steps",
])

# ===== SLIDE 3: PROJECT OVERVIEW =====
add_content_slide("🎯 Project Overview", bullets=[
    "• Fully automated API test generation framework",
    "• Agent-driven workflow — zero manual intervention",
    "• Requirement → Test Cases → Code → Execution → PR (single pipeline)",
    "• Covers 16 REST APIs across 6 modules",
    "• Data-driven approach: JSON files → HashMap → POJOs → API calls",
    "",
    "Benefits:",
    "• 90% reduction in manual test writing effort",
    "• Consistent test patterns across all modules",
    "• Immediate feedback loop on API changes",
    "• Full traceability from requirement to test result",
])

# ===== SLIDE 4: TECH STACK =====
add_content_slide("🛠️ Technology Stack", table_data=[
    ["Category", "Technology", "Version", "Purpose"],
    ["Language", "Java", "11", "Core programming language"],
    ["Build Tool", "Apache Maven", "3.x", "Dependency & build management"],
    ["Test Framework", "TestNG", "7.9.0", "Test execution & reporting"],
    ["API Client", "RestAssured", "5.4.0", "HTTP request/response handling"],
    ["Serialization", "Gson", "2.10.1", "POJO → JSON conversion"],
    ["JSON Parsing", "Jackson Databind", "2.16.1", "File → HashMap deserialization"],
    ["Excel Handling", "Apache POI", "5.2.5", "Test case management"],
    ["Reporting", "Extent Reports", "5.1.1", "Rich HTML test reports"],
    ["Logging", "Log4j 2", "2.22.1", "Application logging"],
])

# ===== SLIDE 5: ARCHITECTURE =====
add_content_slide("🏗️ Framework Architecture", code_text=
"""
    ┌─────────────────────────────────────────────────────────────────────┐
    │                      QA AUTOMATION FRAMEWORK                         │
    ├─────────────────────────────────────────────────────────────────────┤
    │                                                                      │
    │  TEST LAYER: AuthTests | JobsTests | AppTests | ProfileTests        │
    │                              │                                       │
    │                              ▼                                       │
    │  BASE TEST: Authentication + Config Loading (BaseTest.java)         │
    │                              │                                       │
    │                              ▼                                       │
    │  ORCHESTRATION: ActorHelper (Build Payload + Gson + RestUtils)      │
    │                              │                                       │
    │              ┌───────────────┼───────────────┐                      │
    │              ▼               ▼               ▼                      │
    │  POJO Layer        RestUtils (HTTP)     CommonMethod (Data)         │
    │  createPayload()   post/get/put/del     readTestData()              │
    │                              │                                       │
    │                              ▼                                       │
    │  EXTERNAL: HiringManagementSystem REST APIs (16 endpoints)          │
    └─────────────────────────────────────────────────────────────────────┘
""")

# ===== SLIDE 6: DATA FLOW =====
add_content_slide("🔄 Data Flow Pipeline", code_text=
"""
    ┌─────────────┐     ┌───────────────┐     ┌────────────────┐     ┌──────────────┐     ┌─────────┐
    │  JSON File  │────▶│ CommonMethod  │────▶│  ActorHelper   │────▶│  RestUtils   │────▶│   API   │
    │  (testdata) │     │ readTestData()│     │  createJob()   │     │   post()     │     │ Server  │
    │             │     │               │     │                │     │              │     │         │
    │ key:value   │     │ HashMap<S,S>  │     │ POJO→Gson→JSON │     │ HTTP Request │     │ REST    │
    └─────────────┘     └───────────────┘     └────────────────┘     └──────────────┘     └────┬────┘
                                                                                                │
    ┌─────────────┐     ┌───────────────┐     ┌────────────────┐     ┌──────────────┐          │
    │  Test Case  │◀────│  Assertions   │◀────│  ActorHelper   │◀────│  RestUtils   │◀─────────┘
    │  (verify)   │     │ Assert.equals │     │  return resp   │     │  Response    │
    └─────────────┘     └───────────────┘     └────────────────┘     └──────────────┘


    Flow: JSON → HashMap → POJO → Gson → JSON String → HTTP → Response → Assert
""")

# ===== SLIDE 7: AGENT PIPELINE =====
add_content_slide("🤖 Agent-Driven QA Pipeline", table_data=[
    ["Step", "Skill", "Input", "Output", "Quality Gate"],
    ["1", "requirement-understanding", "User Story (.docx)", "Requirements list", "All reqs extracted"],
    ["2", "testcase-generation", "Requirements", "Excel test cases", "Cases documented"],
    ["3", "pojo-evaluation", "API contracts", "Java POJO classes", "createPayload() exists"],
    ["4", "code-generation", "Test cases + POJOs", "TestNG classes", "Compiles without errors"],
    ["5", "test-execution", "Test classes", "Pass/Fail report", "100% pass rate"],
    ["6", "git-push", "Passing code", "Pull Request", "PR created"],
])

# ===== SLIDE 8: API COVERAGE =====
add_content_slide("📊 API Coverage — 16 APIs", table_data=[
    ["Module", "Endpoints", "Methods", "Auth Required"],
    ["🔐 Authentication", "2", "POST", "No"],
    ["💼 Jobs", "5", "GET, POST, PUT, DELETE", "Yes (Recruiter)"],
    ["📄 Applications", "4", "GET, POST, PUT", "Yes (Candidate/Recruiter)"],
    ["👤 User Profile", "2", "GET, PUT", "Yes"],
    ["📂 File Upload", "1", "POST (multipart)", "Yes"],
    ["🛡️ Admin", "2", "GET, DELETE", "Yes (Admin)"],
    ["TOTAL", "16", "Full CRUD", "Role-based"],
])

# ===== SLIDE 9: CLASS DIAGRAM =====
add_content_slide("📐 Class Diagram", code_text=
"""
    ┌────────────────────────────┐       ┌────────────────────────────┐
    │         BaseTest           │       │       URLGenerator         │
    ├────────────────────────────┤       ├────────────────────────────┤
    │ # properties: Properties   │       │ + AUTH_REGISTER: String    │
    │ # accessToken: String      │       │ + JOBS: String             │
    │ # restUtils: RestUtils     │       │ + APPLICATIONS: String     │
    ├────────────────────────────┤       │ + USER_PROFILE: String     │
    │ + setUpWithRole(): String  │       │ + ADMIN_USERS: String      │
    │ + generateAccessToken()    │       └────────────────────────────┘
    │ + getAdminToken()          │
    │ + getRecruiterToken()      │
    │ + getCandidateToken()      │
    └────────────┬───────────────┘
                 │ extends
    ┌────────────▼───────────────┐       ┌────────────────────────────┐
    │        SampleTest          │       │        RestUtils           │
    └────────────────────────────┘       ├────────────────────────────┤
                                         │ + post(url, body): Resp    │
    ┌────────────────────────────┐       │ + get(url): Response       │
    │       ActorHelper          │◆─────▶│ + put(url, body): Resp     │
    ├────────────────────────────┤       │ + delete(url): Response    │
    │ + createJob(): Response    │       │ + uploadFile(): Response   │
    │ + applyForJob(): Response  │       └────────────────────────────┘
    │ + getUserProfile(): Resp   │
    │ + deleteUser(): Response   │       ┌────────────────────────────┐
    └────────────────────────────┘──────▶│     POJO Classes           │
                                         │  + createPayload(): POJO   │
                                         └────────────────────────────┘
""")

# ===== SLIDE 10: SEQUENCE DIAGRAM =====
add_content_slide("🔀 Sequence Diagram — Create Job", code_text=
"""
    Test        CommonMethod   ActorHelper   POJO         RestUtils      API
     │               │              │           │             │            │
     │readTestData() │              │           │             │            │
     │──────────────▶│              │           │             │            │
     │  HashMap      │              │           │             │            │
     │◀──────────────│              │           │             │            │
     │               │              │           │             │            │
     │    createJob(testData)       │           │             │            │
     │─────────────────────────────▶│           │             │            │
     │               │              │createPld()│             │            │
     │               │              │──────────▶│             │            │
     │               │              │ POJO obj  │             │            │
     │               │              │◀──────────│             │            │
     │               │              │           │             │            │
     │               │              │ gson.toJson() → JSON    │            │
     │               │              │           │             │            │
     │               │              │   post(url, json)       │            │
     │               │              │────────────────────────▶│            │
     │               │              │           │             │POST /jobs  │
     │               │              │           │             │───────────▶│
     │               │              │           │             │  201       │
     │               │              │           │             │◀───────────│
     │               │              │  Response │             │            │
     │               │              │◀────────────────────────│            │
     │    Response   │              │           │             │            │
     │◀─────────────────────────────│           │             │            │
     │Assert(201) ✓  │              │           │             │            │
""")

# ===== SLIDE 11: DESIGN DECISIONS =====
add_content_slide("💡 Key Design Decisions", table_data=[
    ["Decision", "Rationale", "Benefit"],
    ["ActorHelper as single entry", "Abstracts API complexity", "Tests stay clean"],
    ["BaseTest for auth only", "Separation of concerns", "Clear responsibility"],
    ["HashMap as data carrier", "Decouples JSON from POJOs", "Flexible data source"],
    ["createPayload() in POJO", "Each POJO owns structure", "Self-contained models"],
    ["Gson for serialization", "Fine-grained JSON control", "Custom payloads easy"],
    ["Exception on non-200", "Fail-fast approach", "Issues caught immediately"],
    ["prettyPrint() everywhere", "Full debugging visibility", "Easy troubleshooting"],
    ["Static URLGenerator", "Single source of truth", "Easy maintenance"],
])

# ===== SLIDE 12: SCALABILITY =====
add_content_slide("📈 Scalability & Extensibility", table_data=[
    ["Scenario", "Action Required", "Files Changed"],
    ["Add new API endpoint", "Add constant to URLGenerator", "1 file"],
    ["Add new API module", "Create POJO + ActorHelper method", "2 files"],
    ["Switch environment", "Update base.url in config", "1 file"],
    ["Add new test scenario", "Add JSON file + test method", "2 files"],
    ["Enable parallel execution", "Update testng.xml", "1 file"],
    ["CI/CD integration", "Add mvn clean test to pipeline", "1 config"],
])

# ===== SLIDE 13: SUMMARY =====
slide = add_title_slide(
    "✅ Summary",
    "16 APIs | 6-Layer Architecture | 6 Agent Skills | 100% Autonomous"
)
txBox = slide.shapes.add_textbox(Inches(1.5), Inches(4.5), Inches(5), Inches(2.5))
tf = txBox.text_frame
tf.word_wrap = True
items = ["Complete layered framework", "Agent-driven pipeline (Req → PR)",
         "Full API coverage (16 endpoints)", "Reusable components", "Data-driven testing"]
p = tf.paragraphs[0]
p.text = "Delivered:"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = BLUE_ACCENT
for item in items:
    p = tf.add_paragraph()
    p.text = f"  ✓ {item}"
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE

txBox2 = slide.shapes.add_textbox(Inches(7), Inches(4.5), Inches(5), Inches(2.5))
tf2 = txBox2.text_frame
tf2.word_wrap = True
items2 = ["CI/CD pipeline (GitHub Actions)", "Performance testing",
          "Contract/Schema testing", "Test metrics dashboard", "Multi-environment"]
p = tf2.paragraphs[0]
p.text = "Next Steps:"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = PURPLE
for item in items2:
    p = tf2.add_paragraph()
    p.text = f"  → {item}"
    p.font.size = Pt(16)
    p.font.color.rgb = LIGHT_GRAY

# ===== SLIDE 14: THANK YOU =====
add_title_slide("Thank You! 🎉", "Questions?")

# Save
output_path = os.path.join(r"C:\Users\odeshpande\CongaAutomation\Hackathon\QA-HiringManagementSystem\docs", "QA-Framework-Presentation.pptx")
prs.save(output_path)
print(f"Presentation saved to: {output_path}")

